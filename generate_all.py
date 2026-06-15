"""
generate_all.py  —  ACADEMIQ
Genere 4 documents professionnels :
  1. ACADEMIQ_CahierDesCharges_BDD.pdf   (ReportLab, en-tetes/pieds de page)
  2. ACADEMIQ_Rapport_Final.pdf           (ReportLab, en-tetes/pieds de page)
  3. ACADEMIQ_V3_Presentation.pptx           (python-pptx, 23 diapositives)
  4. ACADEMIQ_V3_Presentation_Slides.pdf     (ReportLab, paysage)

Usage : python generate_all.py
"""

import os, re
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    BaseDocTemplate, PageTemplate, Frame,
    Paragraph, Spacer, Table, TableStyle,
    PageBreak, HRFlowable, Preformatted, NextPageTemplate
)
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Couleurs ──────────────────────────────────────────────────────────────────
BF  = colors.HexColor('#1A3A5C')
BM  = colors.HexColor('#2C5282')
BC  = colors.HexColor('#4A90D9')
BP  = colors.HexColor('#EBF4FF')
GL  = colors.HexColor('#F5F7FA')
GR  = colors.HexColor('#D8E4F0')
GT  = colors.HexColor('#1A1A1A')

W, H = A4
ML, MR, MT, MB = 2.0*cm, 2.0*cm, 2.8*cm, 2.0*cm
TW = W - ML - MR   # largeur utile


# ═══════════════════════════════════════════════════════════════════════════════
#  STYLES
# ═══════════════════════════════════════════════════════════════════════════════
def make_styles():
    s = {}
    b = getSampleStyleSheet()

    def ps(name, **kw):
        s[name] = ParagraphStyle(name, parent=b['Normal'], **kw)

    ps('body',  fontName='Helvetica', fontSize=9.5, leading=14,
       textColor=GT, alignment=TA_JUSTIFY, spaceAfter=5)
    ps('h1',    fontName='Helvetica-Bold', fontSize=18, leading=22,
       textColor=BF, alignment=TA_CENTER, spaceAfter=8, spaceBefore=14)
    ps('h2',    fontName='Helvetica-Bold', fontSize=13, leading=17,
       textColor=BF, alignment=TA_LEFT, spaceAfter=6, spaceBefore=14)
    ps('h3',    fontName='Helvetica-Bold', fontSize=11, leading=14,
       textColor=BM, alignment=TA_LEFT, spaceAfter=5, spaceBefore=10)
    ps('h4',    fontName='Helvetica-Bold', fontSize=9.5, leading=13,
       textColor=BM, alignment=TA_LEFT, spaceAfter=4, spaceBefore=8)
    ps('pre',   fontName='Courier', fontSize=7.5, leading=10,
       textColor=GT, leftIndent=8, rightIndent=8,
       backColor=colors.HexColor('#F0F4F8'), spaceAfter=8, spaceBefore=4)
    ps('bq',    fontName='Helvetica-Oblique', fontSize=9, leading=13,
       textColor=colors.HexColor('#2C4A6E'),
       backColor=colors.HexColor('#EDF3FB'),
       leftIndent=12, rightIndent=8, spaceAfter=7, spaceBefore=4,
       borderPadding=5)
    ps('li',    fontName='Helvetica', fontSize=9.5, leading=14,
       textColor=GT, leftIndent=14, bulletIndent=4, spaceAfter=2)
    ps('caption', fontName='Helvetica-Oblique', fontSize=8, leading=11,
       textColor=BM, alignment=TA_CENTER, spaceAfter=8, spaceBefore=2)
    # cellules tableau
    ps('th',    fontName='Helvetica-Bold', fontSize=8, leading=10,
       textColor=colors.white, alignment=TA_LEFT)
    ps('th_c',  fontName='Helvetica-Bold', fontSize=8, leading=10,
       textColor=colors.white, alignment=TA_CENTER)
    ps('td',    fontName='Helvetica', fontSize=8, leading=10,
       textColor=GT, alignment=TA_LEFT)
    ps('td_c',  fontName='Helvetica', fontSize=8, leading=10,
       textColor=GT, alignment=TA_CENTER)
    return s

ST = make_styles()


# ═══════════════════════════════════════════════════════════════════════════════
#  EN-TETES ET PIEDS DE PAGE
# ═══════════════════════════════════════════════════════════════════════════════
def make_page_decorator(doc_title, right_header):
    def _draw(canvas, doc):
        canvas.saveState()
        hx, hy, hw, hh = ML, H - MT + 0.25*cm, TW, 0.55*cm
        # En-tete bleu
        canvas.setFillColor(BF)
        canvas.rect(hx, hy, hw, hh, fill=1, stroke=0)
        canvas.setFillColor(colors.white)
        canvas.setFont('Helvetica-Bold', 7.5)
        canvas.drawString(hx + 4, hy + 0.17*cm, doc_title)
        canvas.setFont('Helvetica', 7.5)
        canvas.drawRightString(hx + hw - 4, hy + 0.17*cm, right_header)
        canvas.setStrokeColor(BC)
        canvas.setLineWidth(1.0)
        canvas.line(hx, hy - 1, hx + hw, hy - 1)
        # Pied de page
        fy = MB - 0.4*cm
        canvas.setStrokeColor(GR)
        canvas.setLineWidth(0.5)
        canvas.line(ML, fy + 0.35*cm, ML + TW, fy + 0.35*cm)
        canvas.setFont('Helvetica', 7)
        canvas.setFillColor(BM)
        canvas.drawString(ML, fy, 'Methode Merise  .  MariaDB/InnoDB  .  Django ORM')
        canvas.setFont('Helvetica-Bold', 7)
        canvas.drawRightString(ML + TW, fy, f'Page {doc.page}')
        canvas.restoreState()

    def _first(canvas, doc):
        # Pied de page seulement sur couverture
        canvas.saveState()
        fy = MB - 0.4*cm
        canvas.setStrokeColor(GR)
        canvas.setLineWidth(0.3)
        canvas.line(ML, fy + 0.35*cm, ML + TW, fy + 0.35*cm)
        canvas.restoreState()

    return _draw, _first


# ═══════════════════════════════════════════════════════════════════════════════
#  PARSEUR MARKDOWN -> FLOWABLES REPORTLAB
# ═══════════════════════════════════════════════════════════════════════════════
def _escape(t):
    return t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

def _inline(text):
    t = _escape(text)
    t = re.sub(r'`([^`]+)`',
        r'<font name="Courier" size="8" color="#1A3A5C">\1</font>', t)
    t = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', t)
    t = re.sub(r'\*(.+?)\*',     r'<i>\1</i>', t)
    return t

def _build_table(table_lines):
    """
    Construit un ReportLab Table depuis les lignes Markdown.
    Cle : les cellules sont des Paragraph pour permettre le retour a la ligne.
    """
    rows = []
    for ln in table_lines:
        ln = ln.strip()
        if re.match(r'^\|[-| :]+\|$', ln):
            continue
        cells = [c.strip() for c in ln.strip('|').split('|')]
        rows.append(cells)

    if not rows:
        return None

    ncols = max(len(r) for r in rows)
    rows  = [r + [''] * (ncols - len(r)) for r in rows]

    # Calcul proportionnel des largeurs selon longueur max du contenu
    col_maxlen = []
    for ci in range(ncols):
        ml = max(len(rows[ri][ci]) for ri in range(len(rows)))
        col_maxlen.append(max(ml, 2))
    total_chars = sum(col_maxlen)
    usable = TW - 0.3*cm
    # Largeur minimale basee sur le contenu : 0.20cm/char + 0.5cm padding (capee a 3cm)
    # Evite que les codes courts (ex: RG-MSG06, 8 chars) ne soient coupes sur 2 lignes
    col_w = [
        max(usable * (l / total_chars), min(l * 0.20*cm + 0.5*cm, 3.0*cm), 1.0*cm)
        for l in col_maxlen
    ]
    # Reajustement si depassement
    ratio = usable / sum(col_w)
    if ratio < 1.0:
        col_w = [w * ratio for w in col_w]

    # Construction : chaque cellule = Paragraph (wrap automatique)
    data = []
    for ri, row in enumerate(rows):
        rdata = []
        for ci, cell in enumerate(row):
            txt = _inline(cell.strip())
            if ri == 0:
                sty = ST['th_c'] if ncols > 5 else ST['th']
            else:
                stripped = cell.strip()
                # Valeurs tres courtes (codes) -> centrees ; seuil 4 pour eviter
                # de centrer des mots comme ELEVE(5) PARENT(6) Salles(6) Cours(5)
                if len(stripped) <= 4 and stripped not in ('', ' '):
                    sty = ST['td_c']
                else:
                    sty = ST['td']
            rdata.append(Paragraph(txt, sty))
        data.append(rdata)

    tbl = Table(data, colWidths=col_w, repeatRows=1)
    tbl.setStyle(TableStyle([
        ('BACKGROUND',     (0, 0), (-1, 0),  BF),
        ('TEXTCOLOR',      (0, 0), (-1, 0),  colors.white),
        ('FONTNAME',       (0, 0), (-1, 0),  'Helvetica-Bold'),
        ('FONTSIZE',       (0, 0), (-1, 0),  8),
        ('FONTNAME',       (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE',       (0, 1), (-1, -1), 8),
        ('TEXTCOLOR',      (0, 1), (-1, -1), GT),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, GL]),
        ('GRID',           (0, 0), (-1, -1), 0.25, GR),
        ('LINEBELOW',      (0, 0), (-1, 0),  0.8,  BC),
        ('TOPPADDING',     (0, 0), (-1, -1), 3),
        ('BOTTOMPADDING',  (0, 0), (-1, -1), 3),
        ('LEFTPADDING',    (0, 0), (-1, -1), 5),
        ('RIGHTPADDING',   (0, 0), (-1, -1), 5),
        ('VALIGN',         (0, 0), (-1, -1), 'TOP'),
    ]))
    return tbl


def md_to_flowables(md_text):
    flowables = []
    first_h2 = True   # Le premier h2 ne doit pas inserer PageBreak (page 2 vierge)
    lines = md_text.split('\n')
    i = 0
    in_code = False
    code_lines = []
    table_lines = []
    pending_li = []

    def flush_li():
        for item in pending_li:
            flowables.append(Paragraph(f'&bull;  {_inline(item)}', ST['li']))
        pending_li.clear()

    def flush_table():
        if table_lines:
            tbl = _build_table(table_lines)
            if tbl:
                flowables.append(Spacer(1, 3))
                flowables.append(tbl)
                flowables.append(Spacer(1, 3))
            table_lines.clear()

    while i < len(lines):
        ln = lines[i]

        # Bloc de code
        if ln.startswith('```'):
            if in_code:
                in_code = False
                flowables.append(Spacer(1, 3))
                flowables.append(Preformatted('\n'.join(code_lines), ST['pre']))
                flowables.append(Spacer(1, 4))
                code_lines = []
            else:
                flush_li(); flush_table()
                in_code = True
            i += 1
            continue
        if in_code:
            code_lines.append(ln)
            i += 1
            continue

        s = ln.strip()

        # Tableau
        if s.startswith('|'):
            flush_li()
            table_lines.append(ln)
            i += 1
            continue
        else:
            flush_table()

        # Separateur
        if re.match(r'^-{3,}$', s):
            flush_li()
            flowables.append(HRFlowable(width='100%', thickness=0.5,
                color=GR, spaceAfter=5, spaceBefore=3))
            i += 1
            continue

        # Titres
        m = re.match(r'^(#{1,4})\s+(.+)$', s)
        if m:
            flush_li()
            lvl = len(m.group(1))
            txt = _inline(m.group(2).strip())
            if lvl == 2:
                if not first_h2:
                    flowables.append(PageBreak())
                first_h2 = False
            flowables.append(Paragraph(txt, ST[f'h{lvl}']))
            if lvl in (1, 2):
                flowables.append(HRFlowable(
                    width='100%',
                    thickness=0.8 if lvl == 1 else 0.4,
                    color=BF if lvl == 1 else BC,
                    spaceAfter=5, spaceBefore=2))
            i += 1
            continue

        # Liste non ordonnee
        m = re.match(r'^[-*]\s+(.+)$', s)
        if m:
            pending_li.append(m.group(1))
            i += 1
            continue
        else:
            flush_li()

        # Liste ordonnee
        m = re.match(r'^\d+\.\s+(.+)$', s)
        if m:
            flowables.append(Paragraph(f'&bull;  {_inline(m.group(1))}', ST['li']))
            i += 1
            continue

        # Citation
        if s.startswith('> '):
            flowables.append(Paragraph(_inline(s[2:].strip()), ST['bq']))
            i += 1
            continue

        # Ligne vide
        if not s:
            flowables.append(Spacer(1, 5))
            i += 1
            continue

        # Legende tableau
        if re.match(r'^(tableau|figure)\s+\d', s, re.IGNORECASE):
            flowables.append(Paragraph(_inline(s), ST['caption']))
            i += 1
            continue

        # Paragraphe normal
        flowables.append(Paragraph(_inline(s), ST['body']))
        i += 1

    flush_li()
    flush_table()
    return flowables


# ═══════════════════════════════════════════════════════════════════════════════
#  PAGES DE COUVERTURE (dessin direct sur canvas)
# ═══════════════════════════════════════════════════════════════════════════════
def _cover_common(canvas, titre, sous_titre, badges, infos):
    canvas.saveState()
    # Fond bleu haut
    canvas.setFillColor(BF)
    canvas.rect(0, H*0.44, W, H*0.56, fill=1, stroke=0)
    canvas.setFillColor(colors.white)
    canvas.rect(0, 0, W, H*0.44, fill=1, stroke=0)
    canvas.setFillColor(BC)
    canvas.rect(0, H*0.44 - 0.12*cm, W, 0.24*cm, fill=1, stroke=0)

    # Titre principal
    canvas.setFont('Helvetica-Bold', 36)
    canvas.setFillColor(colors.white)
    canvas.drawCentredString(W/2, H*0.83, 'ACADEMIQ')

    # Sous-titres
    canvas.setFont('Helvetica', 12)
    canvas.setFillColor(colors.HexColor('#BDD7F5'))
    for idx, line in enumerate(sous_titre):
        canvas.drawCentredString(W/2, H*(0.77 - idx*0.04), line)

    # Ligne decorative
    canvas.setStrokeColor(BC)
    canvas.setLineWidth(1.5)
    canvas.line(ML+1*cm, H*0.68, W-ML-1*cm, H*0.68)

    # Badges stats
    bw = TW / len(badges)
    for idx, (val, lbl) in enumerate(badges):
        bx = ML + idx*bw + bw/2
        by = H*0.60
        canvas.setFillColor(BC)
        canvas.roundRect(bx-1.3*cm, by-0.2*cm, 2.6*cm, 1.2*cm, 5, fill=1, stroke=0)
        canvas.setFont('Helvetica-Bold', 20)
        canvas.setFillColor(colors.white)
        canvas.drawCentredString(bx, by+0.55*cm, val)
        canvas.setFont('Helvetica', 8)
        canvas.setFillColor(colors.HexColor('#EBF4FF'))
        canvas.drawCentredString(bx, by+0.1*cm, lbl)

    # Informations (partie basse)
    y = H*0.40
    for label, val in infos:
        canvas.setFont('Helvetica-Bold', 9)
        canvas.setFillColor(BM)
        canvas.drawString(ML+1.5*cm, y, f'{label} :')
        canvas.setFont('Helvetica', 9)
        canvas.setFillColor(GT)
        canvas.drawString(ML+7*cm, y, val)
        y -= 0.68*cm

    # Pied
    canvas.setFont('Helvetica-Oblique', 8)
    canvas.setFillColor(colors.HexColor('#888'))
    canvas.drawCentredString(W/2, MB*0.6, titre)
    canvas.setStrokeColor(GR)
    canvas.setLineWidth(0.4)
    canvas.line(ML, MB, ML+TW, MB)
    canvas.restoreState()


def cover_cdc(canvas, doc):
    _cover_common(canvas,
        titre='Django 4.x  .  MariaDB/InnoDB  .  Methode Merise  .  Bootstrap 5',
        sous_titre=[
            'Cahier des Charges & Conception de Base de Donnees',
            'Methode Merise  .  MariaDB/InnoDB  .  Django ORM',
        ],
        badges=[('24','Tables'), ('59','Regles'), ('38','Cles FK'), ('12','UNIQUE')],
        infos=[
            ('Document',   'Cahier des Charges  —  Conception BDD'),
            ('Version',    'V3  —  Juin 2026'),
            ('Auteur',     'MARIKO Lamine   |   CI0223065023'),
            ('Filiere',    'Master 1  —  Genie Informatique'),
            ('Encadreur',  'Dr. ZEZE'),
            ('Institution','Universite Nangui Abrogoua'),
        ])


def cover_rapport(canvas, doc):
    _cover_common(canvas,
        titre='Django 4.x  .  MariaDB/InnoDB  .  Bootstrap 5  .  ReportLab',
        sous_titre=[
            'Rapport de Projet de Fin d\'Etudes',
            'Conception et Implementation d\'un Systeme Centralise',
            'de Gestion des Activites Scolaires',
        ],
        badges=[('24','Tables'), ('59','Regles'), ('7','Roles'), ('17','Modules')],
        infos=[
            ('Etudiant',   'MARIKO Lamine   |   CI0223065023'),
            ('Filiere',    'Master 1  —  Genie Informatique'),
            ('Encadreur',  'Dr. ZEZE'),
            ('Institution','Universite Nangui Abrogoua'),
            ('Annee',      '2025  —  2026'),
        ])


# ═══════════════════════════════════════════════════════════════════════════════
#  CONSTRUCTION PDF
# ═══════════════════════════════════════════════════════════════════════════════
def build_pdf(md_file, pdf_file, doc_title, right_header, cover_fn):
    print(f'  Lecture : {md_file}')
    with open(md_file, encoding='utf-8') as f:
        md_text = f.read()

    _draw, _first = make_page_decorator(doc_title, right_header)

    frame_cover  = Frame(ML, MB, TW, H - 2*MB,     id='cover')
    frame_normal = Frame(ML, MB, TW, H - MT - MB,  id='normal')

    doc = BaseDocTemplate(pdf_file, pagesize=A4,
        leftMargin=ML, rightMargin=MR, topMargin=MT, bottomMargin=MB)
    doc.addPageTemplates([
        PageTemplate(id='Cover',  frames=[frame_cover],  onPage=cover_fn),
        PageTemplate(id='Normal', frames=[frame_normal], onPage=_draw),
    ])

    story = [NextPageTemplate('Normal'), PageBreak()]
    story += md_to_flowables(md_text)

    print(f'  Generation : {pdf_file}')
    doc.build(story)
    size_kb = os.path.getsize(pdf_file) // 1024
    print(f'  OK  ->  {os.path.basename(pdf_file)}  ({size_kb} Ko)')


# ═══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT (23 diapositives)
# ═══════════════════════════════════════════════════════════════════════════════
BF_r = RGBColor(0x1A, 0x3A, 0x5C)
BM_r = RGBColor(0x2C, 0x52, 0x82)
BC_r = RGBColor(0x4A, 0x90, 0xD9)
BP_r = RGBColor(0xEB, 0xF4, 0xFF)
GL_r = RGBColor(0xF5, 0xF7, 0xFA)
GT_r = RGBColor(0x1A, 0x1A, 0x1A)
WT_r = RGBColor(0xFF, 0xFF, 0xFF)
SW, SH = Inches(13.333), Inches(7.5)


def _bg(sld, rgb):
    sld.background.fill.solid()
    sld.background.fill.fore_color.rgb = rgb


def _rect(sld, l, t, w, h, rgb):
    shp = sld.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    shp.line.fill.background()
    return shp


def _txt(sld, text, l, t, w, h, size, rgb,
         bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = sld.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.color.rgb = rgb
    run.font.bold  = bold
    run.font.italic = italic
    return tb


def _div(sld, y, rgb=None):
    if rgb is None:
        rgb = BC_r
    ln = sld.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(0.035))
    ln.fill.solid()
    ln.fill.fore_color.rgb = rgb
    ln.line.fill.background()


def _header(sld, text):
    _rect(sld, 0, 0, 13.333, 1.1, BF_r)
    _txt(sld, text, 0.4, 0.2, 12.5, 0.8, 24, WT_r, bold=True)


def _ppt_tbl(sld, hdrs, rows, l, t, w, h):
    nr = len(rows) + 1
    nc = len(hdrs)
    tbl = sld.shapes.add_table(nr, nc,
        Inches(l), Inches(t), Inches(w), Inches(h)).table
    cw = Inches(w / nc)
    for ci in range(nc):
        tbl.columns[ci].width = cw
    for ci, hdr in enumerate(hdrs):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = BF_r
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.color.rgb = WT_r
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(8.5)
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BP_r if ri % 2 == 0 else WT_r
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.runs[0].font.size = Pt(8.5)
            p.runs[0].font.color.rgb = GT_r


def _bullets(sld, items, x, y, w, size=10.5, color=None):
    if color is None:
        color = GT_r
    for item in items:
        _txt(sld, f'   {item}', x, y, w, 0.42, size, color)
        y += 0.42
    return y


# ── 23 slides ─────────────────────────────────────────────────────────────────

def s_title(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, BF_r)
    _rect(sld, 0, 4.5, 13.333, 3.0, RGBColor(0x0D, 0x1E, 0x30))
    _rect(sld, 0, 4.44, 13.333, 0.12, BC_r)
    _txt(sld, 'ACADEMIQ', 0.5, 1.0, 12.3, 1.4, 52, WT_r, bold=True, align=PP_ALIGN.CENTER)
    _txt(sld, 'Rapport de Projet de Fin d\'Etudes', 0.5, 2.45, 12.3, 0.55, 17,
         RGBColor(0xBD,0xD7,0xF5), align=PP_ALIGN.CENTER)
    _txt(sld, 'Conception & Implementation d\'un Systeme Centralise de Gestion des Activites Scolaires',
         0.5, 3.0, 12.3, 0.9, 12, RGBColor(0x9B,0xC1,0xE8), align=PP_ALIGN.CENTER)
    for idx, (val, lbl) in enumerate([('24','Tables'),('59','Regles'),('7','Roles'),('17','Modules')]):
        bx = 1.0 + idx*3.0
        _rect(sld, bx, 4.7, 2.2, 1.5, BC_r)
        _txt(sld, val, bx, 4.72, 2.2, 0.88, 30, WT_r, bold=True, align=PP_ALIGN.CENTER)
        _txt(sld, lbl, bx, 5.58, 2.2, 0.42, 10, RGBColor(0xEB,0xF4,0xFF), align=PP_ALIGN.CENTER)
    _txt(sld, 'MARIKO Lamine  |  CI0223065023  |  M1 Genie Informatique  |  Univ. Nangui Abrogoua  |  Juin 2026',
         0.5, 6.9, 12.3, 0.4, 8.5, RGBColor(0x7F,0xA8,0xCC), align=PP_ALIGN.CENTER)


def s_plan(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'Sommaire')
    chapters = [
        ('I',   'Introduction et Contexte'),
        ('II',  'Analyse et Specification des Besoins'),
        ('III', 'Conception du Systeme (Methode Merise)'),
        ('IV',  'Implementation et Developpement (Django)'),
        ('V',   'Tests, Validation et Resultats'),
        ('VI',  'Conclusion et Perspectives'),
    ]
    y = 1.35
    for num, title in chapters:
        _rect(sld, 0.4, y, 1.0, 0.52, BF_r)
        _txt(sld, num, 0.4, y+0.07, 1.0, 0.4, 13, WT_r, bold=True, align=PP_ALIGN.CENTER)
        _txt(sld, title, 1.55, y+0.08, 11.0, 0.42, 13, GT_r)
        y += 0.82


def s_contexte(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'I. Contexte et Problematique')
    _txt(sld, 'Situation actuelle', 0.4, 1.2, 5.9, 0.4, 13, BF_r, bold=True)
    _txt(sld, 'Objectif du projet',  7.0, 1.2, 6.0, 0.4, 13, BF_r, bold=True)
    _div(sld, 1.62)
    probs = [
        '•  Notes dans des cahiers de texte individuels',
        '•  Absences sur fiches papier circulantes',
        '•  Bulletins produits manuellement',
        '•  Distribution 2 a 3 semaines apres la cloture',
        '•  Aucune traçabilite des modifications',
        '•  Communication orale non structuree',
    ]
    objs = [
        '✔  Centraliser toutes les donnees',
        '✔  Calculer automatiquement moyennes et rangs',
        '✔  Generer les bulletins PDF en 1 clic',
        '✔  Notifier les parents en temps reel',
        '✔  Garantir la traçabilite de chaque action',
        '✔  Securiser les acces par 7 roles cloisonnes',
    ]
    y = 1.73
    for p, o in zip(probs, objs):
        _txt(sld, p, 0.4, y, 6.0, 0.4, 10, RGBColor(0xBB,0x33,0x33))
        _txt(sld, o, 7.0, y, 6.0, 0.4, 10, RGBColor(0x22,0x7A,0x33))
        y += 0.44


def s_acteurs(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'II. Les 7 Roles Utilisateurs')
    roles = [
        (BF_r,                         'DIRECTION',      'Acces total : gestion comptes, annees scolaires, historique'),
        (BM_r,                         'ADMINISTRATION', 'Cours, emplois du temps, salles, matieres'),
        (BC_r,                         'SCOLARITE',      'Inscriptions, absences, generation et export bulletins'),
        (RGBColor(0x1A,0x6B,0x4A),    'FINANCES',       'Tarifs, paiements, export bilan financier PDF'),
        (RGBColor(0x6B,0x3A,0x1A),    'ENSEIGNANT',     'Notes et absences pour ses propres cours uniquement'),
        (RGBColor(0x4A,0x1A,0x6B),    'ELEVE',          'Consultation de son propre dossier (notes, bulletins, EDT)'),
        (RGBColor(0x1A,0x4A,0x6B),    'PARENT',         'Consultation du dossier de ses enfants, messagerie'),
    ]
    y = 1.25
    for col, role, desc in roles:
        _rect(sld, 0.4, y, 2.8, 0.5, col)
        _txt(sld, role, 0.4, y+0.08, 2.8, 0.38, 10, WT_r, bold=True, align=PP_ALIGN.CENTER)
        _txt(sld, desc, 3.35, y+0.09, 9.5, 0.38, 10, GT_r)
        y += 0.58


def s_permissions(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'II. Matrice des Permissions (extrait)')
    _txt(sld, 'C=Creer  L=Lire  M=Modifier  S=Supprimer  —=Aucun',
         0.4, 1.12, 12.5, 0.35, 9, BM_r, italic=True)
    hdrs = ['Fonction', 'DIR', 'ADM', 'SCO', 'FIN', 'ENS', 'ELEV', 'PAR']
    data = [
        ('Annees scolaires',  'CLMS', '—',    'L',   'L',   'L',        'L',         'L'),
        ('Cours',             'CLMS', 'CLMS', 'L',   'L',   'L(siens)', '—',         '—'),
        ('Saisie notes',      '—',    '—',    '—',   '—',   'CLM',      '—',         '—'),
        ('Consultation notes','LM',   'L',    'L',   '—',   'L(siens)', 'L(siennes)','L(enf.)'),
        ('Absences',          'CLMS', '—',    'CLM', '—',   'CLM',      'L',         'L(enf.)'),
        ('Bulletins PDF',     'CL',   'L',    'CL',  '—',   'L(elev)',  'L(sien)',   'L(enf.)'),
        ('Paiements',         'CLMS', '—',    '—',   'CLMS','—',        '—',         'L(enf.)'),
        ('Messagerie',        'CLM',  'CLM',  'CLM', 'CLM', 'CLM(tous)','CLM(pers)', 'CLM(pers)'),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 1.55, 12.7, 5.65)


def s_regles(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'II. Regles de Gestion — Extrait (59 au total)')
    hdrs = ['Code', 'Domaine', 'Regle de gestion']
    data = [
        ('RG-T02',  'Temporel',   'Une seule ANNEE_SCOLAIRE active a la fois (signal pre_save Django)'),
        ('RG-C03',  'Classe',     'Un ELEVE dans 1 seule CLASSE par ANNEE — UNIQUE(eleve_id, annee_id)'),
        ('RG-N03',  'Note',       'Valeur NOTE obligatoirement dans [0 ; 20] — contrainte CHECK SQL'),
        ('RG-N06',  'Note',       'Saisie bloquee si PERIODE.cloturee = True (verification dans la vue)'),
        ('RG-M06',  'Moyenne',    'RESULTAT_MATIERE alimente par signaux uniquement — aucune saisie directe'),
        ('RG-AB02', 'Absence',    "Validation des absences : competence de l'ENSEIGNANT titulaire"),
        ('RG-AB06', 'Absence',    'Seuil 20h absences non justifiees : NOTIFICATION automatique envoyee'),
        ('RG-B01',  'Bulletin',   'Un seul BULLETIN par (ELEVE, PERIODE) — UNIQUE constraint en base'),
        ('RG-B02',  'Bulletin',   'Generation possible uniquement apres cloture de la PERIODE'),
        ('RG-FIN09','Finance',    'Export PDF bilan financier accessible a FINANCES et DIRECTION'),
        ('RG-MSG03','Messagerie', "ENSEIGNANT peut ecrire a tout utilisateur actif (y compris ELEVE)"),
        ('RG-A03',  'Securite',   "Acces non autorise : message d'erreur explicite + redirect dashboard"),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 1.2, 12.7, 6.05)


def s_archi(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'III. Architecture Trois Couches')
    layers = [
        (BC_r,  1.25, 'COUCHE PRESENTATION',
         'Navigateur web   HTML5 + CSS3   Bootstrap 5 + Chart.js   Rendu cote serveur (pas de SPA)'),
        (BM_r,  3.0,  'COUCHE APPLICATIVE',
         'Django 4.x (Python)   Routing . Vues FBV . Templates . Signaux . Permissions . ReportLab'),
        (BF_r,  4.75, 'COUCHE DONNEES',
         'MariaDB 10.x InnoDB   24 tables . 95 attributs . 38 cles FK . Transactions ACID'),
    ]
    arrows = [(4.42, 'HTTP / HTTPS'), (4.42 + 1.75, 'ORM Django (SQL)')]
    for col, ty, titre, sub in layers:
        _rect(sld, 0.5, ty, 12.3, 1.45, col)
        _txt(sld, titre, 0.7, ty+0.1, 12.0, 0.5, 13, WT_r, bold=True)
        _txt(sld, sub,   0.7, ty+0.62, 12.0, 0.55, 9.5, RGBColor(0xC5,0xD8,0xEC))
    for ay, albl in arrows:
        _rect(sld, 5.9, ay, 1.5, 0.28, RGBColor(0xCC,0xCC,0xCC))
        _txt(sld, albl, 5.9, ay+0.04, 1.5, 0.22, 7.5, GT_r, align=PP_ALIGN.CENTER)


def s_merise(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'III. Methode Merise — Demarche Adoptee')
    steps = [
        ('MCD', 'Modele Conceptuel de Donnees',
         'Entites, associations, cardinalites.\nIndependant de tout SGBD.'),
        ('MLD', 'Modele Logique de Donnees',
         'Tables, cles primaires et etrangeres.\nPassage Merise formalise.'),
        ('MPD', 'Modele Physique de Donnees',
         'Script SQL MariaDB/InnoDB.\nTypes, contraintes, index, UNIQUE.'),
    ]
    for idx, (code, titre, desc) in enumerate(steps):
        xp = 0.5 + idx*4.3
        _rect(sld, xp, 1.3, 3.8, 0.95, BF_r)
        _txt(sld, code, xp, 1.35, 3.8, 0.55, 22, WT_r, bold=True, align=PP_ALIGN.CENTER)
        _rect(sld, xp, 2.25, 3.8, 2.8, BP_r)
        _txt(sld, titre, xp+0.1, 2.35, 3.6, 0.5, 11, BF_r, bold=True)
        _txt(sld, desc,  xp+0.1, 2.9,  3.6, 1.9, 10, GT_r)
    for xp in [4.35, 8.65]:
        _rect(sld, xp, 1.62, 0.5, 0.28, BC_r)
        _txt(sld, '>>>', xp, 1.64, 0.5, 0.25, 10, WT_r, align=PP_ALIGN.CENTER)

    _txt(sld, 'Justification du choix', 0.4, 5.25, 12.5, 0.4, 12, BF_r, bold=True)
    _div(sld, 5.65)
    for r in [
        '✔  Rigueur : demarche structuree garantissant la coherence avant toute implementation',
        '✔  Independance technologique : le MLD est produit independamment du SGBD cible',
        '✔  Normalisation : la demarche Merise conduit naturellement a un schema en 3FN',
    ]:
        _txt(sld, r, 0.5, 5.7, 12.3, 0.4, 10, RGBColor(0x22,0x7A,0x33))
        pass
    y = 5.72
    for r in [
        '✔  Rigueur : demarche structuree garantissant la coherence avant toute implementation',
        '✔  Independance technologique : le MLD est produit independamment du SGBD cible',
        '✔  Normalisation : la demarche Merise conduit naturellement a un schema en 3FN',
    ]:
        _txt(sld, r, 0.5, y, 12.3, 0.4, 10, RGBColor(0x22,0x7A,0x33))
        y += 0.43


def s_entites(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'III. Les 24 Entites — 6 Domaines Fonctionnels')
    domains = [
        ('Referentiel',  BF_r,  ['PERSONNE', 'PERSONNEL', 'ENSEIGNANT', 'ELEVE', 'PARENT', 'LIEN_PE']),
        ('Pedagogique',  BM_r,  ['ANNEE_SCOLAIRE', 'PERIODE', 'CLASSE', 'SALLE', 'MATIERE', 'COURS', 'EDT']),
        ('Scolarite',    BC_r,  ['INSCRIPTION']),
        ('Evaluation',   RGBColor(0x1A,0x6B,0x4A), ['NOTE', 'RESULTAT_MATIERE', 'BULLETIN', 'ABSENCE']),
        ('Finances',     RGBColor(0x8B,0x4A,0x1A), ['TARIF_NIVEAU', 'FRAIS_SCOLARITE', 'PAIEMENT']),
        ('Communication',RGBColor(0x4A,0x1A,0x8B), ['MESSAGE', 'NOTIFICATION', 'EVENEMENT', 'HISTORIQUE']),
    ]
    pos = [(0.25,1.2),(4.7,1.2),(9.15,1.2),(0.25,3.8),(4.7,3.8),(9.15,3.8)]
    for (dom, col, tbls), (lx, ly) in zip(domains, pos):
        _rect(sld, lx, ly, 4.0, 0.45, col)
        _txt(sld, dom, lx, ly+0.06, 4.0, 0.35, 11, WT_r, bold=True, align=PP_ALIGN.CENTER)
        y2 = ly+0.47
        for tbl in tbls:
            _rect(sld, lx, y2, 4.0, 0.37, BP_r)
            _txt(sld, tbl, lx+0.15, y2+0.05, 3.8, 0.28, 9, BF_r)
            y2 += 0.38


def s_mld(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'III. MLD — Contraintes Cles (38 FK, 12 UNIQUE composites)')
    hdrs = ['Table', 'Cles et contraintes', 'Politique']
    data = [
        ('COURS',           'UNIQUE(matiere_id, classe_id, annee_id)',          'RESTRICT'),
        ('INSCRIPTION',     'UNIQUE(eleve_id, annee_id)',                       'RESTRICT'),
        ('NOTE',            'CHECK(valeur >= 0 AND valeur <= 20)',              '—'),
        ('RESULTAT_MAT.',   'UNIQUE(eleve_id, cours_id, periode_id)',           'RESTRICT'),
        ('BULLETIN',        'UNIQUE(eleve_id, periode_id)',                     'RESTRICT'),
        ('EMPLOI DU TEMPS', 'UNIQUE(salle_id, periode_id, jour, heure_debut)', 'RESTRICT'),
        ('FRAIS_SCOL.',     'UNIQUE(eleve_id, annee_id)',                       'RESTRICT'),
        ('PAIEMENT',        'recu_numero UNIQUE',                               'RESTRICT'),
        ('TARIF_NIVEAU',    'UNIQUE(niveau, annee_id)',                         'RESTRICT'),
        ('HISTORIQUE',      'auteur_id FK -> PERSONNE',                        'SET NULL'),
        ('NOTIFICATION',    'classe_id FK -> CLASSE',                          'SET NULL'),
        ('LIEN_PE',         'UNIQUE(parent_id, eleve_id)',                      'RESTRICT'),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 1.2, 12.7, 6.1)


def s_user_model(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'IV. Modele Utilisateur Personnalise (AbstractBaseUser)')
    _txt(sld, 'Pourquoi AbstractBaseUser ?', 0.4, 1.2, 12.5, 0.4, 13, BF_r, bold=True)
    _div(sld, 1.62)
    items = [
        '•  Email comme identifiant de connexion (USERNAME_FIELD = "email")',
        '•  Tous les acteurs dans une seule table PERSONNE — simplifie les jointures',
        '•  Champs supplementaires : nom, prenom, photo_profil, actif directement',
        '•  Compatibilite parfaite avec le modele Merise (PERSONNE = super-entite)',
        '•  Choix irreversible — doit etre fait avant la premiere migration Django',
    ]
    y = 1.7
    for item in items:
        _txt(sld, item, 0.5, y, 12.3, 0.4, 10.5, GT_r)
        y += 0.44
    _txt(sld, 'Configuration :   AUTH_USER_MODEL = "core.Personne"', 0.4, 4.1, 12.5, 0.4, 11, BF_r, bold=True)
    _div(sld, 4.5)
    hdrs = ['Sous-entite', 'Champs specifiques', 'Groupe Django associe']
    data = [
        ('PERSONNEL',  'fonction (ENUM), date_embauche',   'DIRECTION, ADMINISTRATION, SCOLARITE, FINANCES'),
        ('ENSEIGNANT', 'specialite, grade, date_embauche', 'ENSEIGNANT'),
        ('ELEVE',      'matricule (UNIQUE), date_naiss.',  'ELEVE'),
        ('PARENT',     'telephone, profession',            'PARENT'),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 4.6, 12.7, 2.1)


def s_signaux(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'IV. Signaux Django — Automatisation des Traitements Metier')
    sigs = [
        (BF_r,  'post_save sur NOTE',
         ['Declenche a chaque creation/modification de NOTE',
          'Recalcule la moyenne de l\'eleve pour le cours + periode',
          'Trie toutes les moyennes du cours par ordre decroissant',
          'Attribue les rangs avec bulk_update(["rang"])']),
        (BM_r,  'pre_save sur ANNEE_SCOLAIRE',
         ['Intercepte toute activation d\'une annee scolaire',
          'Desactive toutes les autres annees en une seule requete',
          'Garantit l\'unicite de l\'annee active (RG-T02)',
          'Transparent pour la vue — zero logique en vue']),
        (BC_r,  'post_save sur BULLETIN',
         ['Envoie une NOTIFICATION a l\'eleve concerne',
          'Envoie une NOTIFICATION a chaque parent lie',
          'type_notif = "bulletin" pour filtrage frontend',
          'Aucun appel explicite necessite dans la vue']),
    ]
    y = 1.2
    for col, titre, items in sigs:
        _rect(sld, 0.3, y, 12.7, 0.4, col)
        _txt(sld, titre, 0.45, y+0.07, 12.4, 0.3, 11, WT_r, bold=True)
        y += 0.42
        for item in items:
            _txt(sld, f'    •  {item}', 0.3, y, 12.5, 0.36, 9.5, GT_r)
            y += 0.36
        y += 0.18


def s_bug(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'IV. Bug Identifie & Corrige : bulk_create et signaux Django')
    _txt(sld, 'Probleme', 0.4, 1.2, 5.9, 0.4, 13, RGBColor(0xCC,0x33,0x33), bold=True)
    _txt(sld, 'Solution', 7.0, 1.2, 6.0, 0.4, 13, RGBColor(0x22,0x7A,0x33), bold=True)
    _div(sld, 1.62)
    probs = [
        '•  populate_data utilise bulk_create pour la performance',
        '•  bulk_create() ne declenche PAS le signal post_save',
        '•  Le recalcul des rangs est donc ignore lors de l\'import',
        '•  Tous les rangs restent a NULL apres l\'initialisation',
        '•  Les bulletins affichent "rang : —" au lieu du numero',
    ]
    sols = [
        '•  groupby(resultats, key=lambda r: r.cours_id)',
        '•  enumerate(group, start=1) => res.rang = rang',
        '•  ResultatMatiere.objects.bulk_update(lst, ["rang"])',
        '•  Appele explicitement apres chaque bulk_create',
        '•  Resultat : 840 / 840 rangs calcules, 0 valeur NULL',
    ]
    y = 1.72
    for p, s in zip(probs, sols):
        _txt(sld, p, 0.4, y, 5.9, 0.38, 10, RGBColor(0xBB,0x33,0x33))
        _txt(sld, s, 7.0, y, 6.0, 0.38, 10, RGBColor(0x22,0x7A,0x33))
        y += 0.42
    _rect(sld, 0.3, 4.2, 12.7, 0.07, BC_r)
    _txt(sld, 'Leçon : toujours recalculer les agregats apres un bulk_create (signals non declenches)',
         0.4, 4.35, 12.5, 0.5, 10.5, BF_r, italic=True)


def s_acces(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'IV. Controle d\'Acces a Deux Niveaux')
    _txt(sld, 'Double cloisonnement : Vue + Queryset', 0.4, 1.2, 12.5, 0.4, 13, BF_r, bold=True)
    _div(sld, 1.62)
    hdrs = ['Niveau', 'Mecanisme', 'Exemple concret']
    data = [
        ('Vue',
         '@role_required("ENSEIGNANT")',
         'Seuls les membres du groupe ENSEIGNANT peuvent atteindre /enseignant/notes/'),
        ('Queryset',
         'Cours.objects.filter(enseignant=request.user)',
         "L'enseignant ne voit que ses propres cours dans la liste"),
        ('Objet',
         'get_object_or_404(Cours, pk=id, enseignant=user)',
         "get_object_or_404 retourne 404 si l'objet appartient a un autre"),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 1.7, 12.7, 2.4)

    _txt(sld, 'Redirection intelligente en cas d\'acces refuse', 0.4, 4.3, 12.5, 0.4, 13, BF_r, bold=True)
    _div(sld, 4.72)
    _txt(sld,
         'Au lieu d\'une page HTTP 403 generique (incomprehensible pour l\'utilisateur), '
         'tout acces non autorise affiche un message d\'erreur explicite ET redirige vers '
         'le dashboard propre au role de l\'utilisateur connecte (RG-A03).',
         0.4, 4.78, 12.5, 0.75, 10.5, GT_r)
    _rect(sld, 0.3, 5.62, 12.7, 0.07, BC_r)
    _txt(sld, 'Aucune information sur l\'existence de la page n\'est divulguee a l\'utilisateur non autorise.',
         0.4, 5.75, 12.5, 0.4, 10, BM_r, italic=True)


def s_msg(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'IV. Module Messagerie — Recherche Dynamique')
    _txt(sld, 'Probleme initial', 0.4, 1.2, 5.9, 0.4, 13, RGBColor(0xCC,0x33,0x33), bold=True)
    _txt(sld, 'Solution retenue', 7.0, 1.2, 6.0, 0.4, 13, RGBColor(0x22,0x7A,0x33), bold=True)
    _div(sld, 1.62)
    _txt(sld,
         'Avec des centaines d\'utilisateurs, un simple <select> deroulant est impraticable.\n'
         'La technique style.display="none" sur <option> est non fiable selon les navigateurs.',
         0.4, 1.7, 5.9, 0.85, 10, GT_r)
    _txt(sld,
         'Injection d\'un tableau JSON DEST_DATA dans la page.\n'
         'Reconstruction dynamique du <select> via JavaScript a chaque frappe.\n'
         'Aucun rechargement de page. Compatible tous navigateurs.',
         7.0, 1.7, 6.0, 0.85, 10, GT_r)

    _txt(sld, 'Filtrage des destinataires par role (RG-MSG03)', 0.4, 2.75, 12.5, 0.4, 12, BF_r, bold=True)
    _div(sld, 3.15)
    hdrs = ['Role expediteur', 'Destinataires autorises dans la liste']
    data = [
        ('DIRECTION / ADMIN / SCOLARITE / FINANCES', 'Tous les utilisateurs actifs du systeme'),
        ('ENSEIGNANT',  'Tous les utilisateurs actifs (y compris ELEVE et PARENT)'),
        ('ELEVE',       'Personnel administratif et ENSEIGNANT uniquement'),
        ('PARENT',      'Personnel administratif et ENSEIGNANT uniquement'),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 3.25, 12.7, 2.15)
    _txt(sld, 'Piece jointe : fichiers PDF uniquement — validation Content-Type cote serveur (Django)',
         0.4, 5.52, 12.5, 0.4, 10, BM_r, italic=True)


def s_pdf(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'IV. Generation PDF des Bulletins (ReportLab)')
    steps = [
        ('1', 'Verification',   'PERIODE.cloturee = True ET PERIODE.date_fin <= date.today()'),
        ('2', 'Recuperation',   'RESULTAT_MATIERE de l\'eleve pour la periode concernee'),
        ('3', 'Moy. generale',  'Calcul : Somme(moy x coeff) / Somme(coefficients)'),
        ('4', 'Rang classe',    'Tri decroissant des moyennes generales de tous bulletins classe'),
        ('5', 'Stockage',       'BULLETIN cree — donnees immutables apres generation (RG-B02)'),
        ('6', 'PDF ReportLab',  'Genere en memoire via io.BytesIO — aucun stockage disque'),
        ('7', 'Notification',   'NOTIFICATION automatique vers l\'eleve ET chaque parent lie'),
        ('8', 'Reponse HTTP',   'Content-Type: application/pdf — telechargement immediat'),
    ]
    y = 1.25
    for num, titre, detail in steps:
        _rect(sld, 0.3, y, 0.55, 0.5, BF_r)
        _txt(sld, num, 0.3, y+0.07, 0.55, 0.38, 13, WT_r, bold=True, align=PP_ALIGN.CENTER)
        _txt(sld, titre,  1.0, y+0.08, 2.8, 0.38, 10, BM_r, bold=True)
        _txt(sld, detail, 3.95, y+0.08, 9.0, 0.38, 10, GT_r)
        y += 0.58


def s_tests(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'V. Tests Fonctionnels — Resultats')
    hdrs = ['Scenario de test', 'Resultat attendu', 'Statut']
    data = [
        ('Login email + MDP valides',              'Redirection vers le bon dashboard selon role',      'OK'),
        ('Login avec compte desactive',             'Refus + message d\'erreur clair',                  'OK'),
        ('Acces URL d\'un autre role',              'Message d\'erreur + redirect dashboard propre',    'OK'),
        ('Saisie note valide (ex : 14.5)',          'Note enregistree, moyenne et rangs recalcules',    'OK'),
        ('Saisie note hors [0 ; 20]',              'Formulaire rejete, message d\'erreur affiché',     'OK'),
        ('Saisie note sur periode cloturee',        'Saisie bloquee avec message explicite',            'OK'),
        ('Absence par enseignant titulaire',        'Statut en_attente, notification eleve/parent',     'OK'),
        ('Validation absence par non-titulaire',    'Acces refuse (get_object_or_404 filtre)',          'OK'),
        ('Generation bulletin sans cloture',        'Action bloquee (verification avant traitement)',   'OK'),
        ('Export bilan financier par ELEVE',        'Acces refuse, redirect dashboard ELEVE',           'OK'),
        ('Messagerie ELEVE vers autre ELEVE',       'Destinataire absent de la liste (filtre role)',    'OK'),
        ('Enregistrement paiement par FINANCES',    'Frais mis a jour, statut recalcule automatiquement','OK'),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 1.2, 12.7, 6.05)


def s_integrite(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'V. Tests d\'Integrite — Apres populate_data')
    _txt(sld, 'Commande : python manage.py populate_data', 0.4, 1.15, 12.5, 0.35, 10, BM_r, italic=True)
    hdrs = ['Verification en base', 'Valeur attendue', 'Valeur obtenue', 'Statut']
    data = [
        ('Nombre de RESULTAT_MATIERE',           '840',       '840',       'OK'),
        ('Avec rang != NULL (calcule)',           '840',       '840',       'OK'),
        ('Avec rang NULL (bug precedent)',        '0',         '0',         'OK'),
        ('Doublons BULLETIN(eleve, periode)',     '0',         '0',         'OK'),
        ('Doublons numeros de reçu',             '0',         '0',         'OK'),
        ('Violation CHECK valeur [0;20]',        'Impossible','Impossible', 'OK'),
        ('Doublon INSCRIPTION(eleve, annee)',     '0',         '0',         'OK'),
        ('Annees actives simultanement',         '1',         '1',         'OK'),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 1.6, 12.7, 5.55)


def s_bilan(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'V. Bilan Quantitatif du Projet')
    kpis = [
        ('24',  'Tables BDD'),    ('95',  'Attributs'),
        ('38',  'Cles FK'),       ('12',  'UNIQUE comp.'),
        ('59',  'Regles metier'), ('7',   'Roles users'),
        ('6',   'Apps Django'),   ('17',  'Modules'),
        ('840', 'Resultats/rangs'),
    ]
    y = 1.25
    x_start = 0.3
    col_w = 4.3
    for idx, (val, lbl) in enumerate(kpis):
        col = idx % 3
        row = idx // 3
        lx = x_start + col * col_w
        ly = y + row * 1.3
        _rect(sld, lx, ly, 4.0, 1.12, BF_r)
        _txt(sld, val, lx, ly+0.08, 4.0, 0.65, 26, WT_r, bold=True, align=PP_ALIGN.CENTER)
        _txt(sld, lbl, lx, ly+0.72, 4.0, 0.35, 9,  RGBColor(0xEB,0xF4,0xFF), align=PP_ALIGN.CENTER)


def s_difficultes(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'VI. Difficultes Rencontrees & Solutions Apportees')
    hdrs = ['Difficulte', 'Cause identifiee', 'Solution adoptee']
    data = [
        ('Rangs NULL apres import',
         'bulk_create() ne declenche pas post_save',
         'Appel explicite de bulk_update(["rang"]) avec groupby apres chaque bulk_create'),
        ('Select destinataire bug',
         'style.display="none" sur <option> non fiable',
         'Reconstruction dynamique du DOM via JS (tableau DEST_DATA JSON)'),
        ('HTTP 403 brut pour utilisateurs',
         'role_required redirige vers page 403 generique',
         '_redirect_to_dashboard() avec message d\'erreur explicite par role'),
        ('Mauvaise sidebar pour ENSEIGNANT',
         'Template heritait de base_personnel.html',
         'Template dedie base_enseignant.html cree pour les annonces'),
        ('Liens dashboard inactifs',
         'Cartes accessibles sans filtrage de role',
         'Verifications {% if "ROLE" in groupes_user %} dans les templates'),
    ]
    _ppt_tbl(sld, hdrs, data, 0.3, 1.2, 12.7, 4.2)
    _txt(sld, 'Leçons apprises', 0.4, 5.55, 12.5, 0.4, 12, BF_r, bold=True)
    _div(sld, 5.95)
    _txt(sld,
         '✔  Le cloisonnement doit etre double : niveau VUE et niveau QUERYSET\n'
         '✔  Tester bulk_create separement car les signaux Django ne sont pas declenches',
         0.4, 6.0, 12.5, 0.75, 10, RGBColor(0x22,0x7A,0x33))


def s_conclusion(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, WT_r)
    _header(sld, 'VI. Conclusion & Perspectives')
    _txt(sld, 'Bilan du projet', 0.4, 1.2, 12.5, 0.4, 13, BF_r, bold=True)
    _div(sld, 1.62)
    bilans = [
        '✔  Systeme d\'information scolaire complet et operationnel — 17 modules livres',
        '✔  Architecture Django/MariaDB robuste en 3FN — 38 FK, 12 contraintes UNIQUE',
        '✔  Controle d\'acces strict par 7 roles — aucun acces croise possible en test',
        '✔  Automatisation totale : moyennes, rangs, bulletins, notifications',
    ]
    y = 1.72
    for b in bilans:
        _txt(sld, b, 0.5, y, 12.3, 0.38, 10.5, RGBColor(0x22,0x7A,0x33))
        y += 0.42

    _txt(sld, 'Court terme (3-6 mois)', 0.4, 3.5, 6.0, 0.4, 12, BF_r, bold=True)
    _txt(sld, 'Long terme (> 18 mois)',  7.0, 3.5, 6.0, 0.4, 12, BF_r, bold=True)
    _div(sld, 3.9)
    ct = [
        '•  Import CSV/Excel listes et notes historiques',
        '•  Tests unitaires automatises (pytest-django)',
        '•  Audit accessibilite WCAG 2.1',
    ]
    lt = [
        '•  API REST + application mobile iOS/Android',
        '•  Notifications push temps reel (WebSockets)',
        '•  Architecture multi-etablissements (cloud)',
    ]
    y = 4.0
    for c, l in zip(ct, lt):
        _txt(sld, c, 0.4, y, 6.0, 0.4, 10.5, GT_r)
        _txt(sld, l, 7.0, y, 6.0, 0.4, 10.5, GT_r)
        y += 0.44


def s_questions(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sld, BF_r)
    _rect(sld, 0, 4.4, 13.333, 3.1, RGBColor(0x0D,0x1E,0x30))
    _rect(sld, 0, 4.34, 13.333, 0.12, BC_r)
    _txt(sld, 'Merci pour votre attention', 0.5, 1.4, 12.3, 1.2, 34,
         WT_r, bold=True, align=PP_ALIGN.CENTER)
    _txt(sld, 'ACADEMIQ', 0.5, 3.0, 12.3, 0.7, 22,
         RGBColor(0xBD,0xD7,0xF5), bold=True, align=PP_ALIGN.CENTER)
    _div(sld, 3.75, BC_r)
    _txt(sld, 'Questions & Discussion', 0.5, 3.85, 12.3, 0.6, 16,
         RGBColor(0x9B,0xC1,0xE8), align=PP_ALIGN.CENTER)
    _txt(sld, 'MARIKO Lamine  |  CI0223065023  |  M1 Genie Informatique', 0.5, 5.0, 12.3, 0.4, 11,
         RGBColor(0x7F,0xA8,0xCC), align=PP_ALIGN.CENTER)
    _txt(sld, 'Universite Nangui Abrogoua  |  Juin 2026  |  Encadreur : Dr. ZEZE', 0.5, 5.5, 12.3, 0.4, 11,
         RGBColor(0x7F,0xA8,0xCC), align=PP_ALIGN.CENTER)


def build_pptx(pptx_file):
    print(f'  Generation PPTX : {pptx_file}')
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    slides_fns = [
        s_title, s_plan, s_contexte, s_acteurs, s_permissions, s_regles,
        s_archi, s_merise, s_entites, s_mld,
        s_user_model, s_signaux, s_bug, s_acces, s_msg, s_pdf,
        s_tests, s_integrite, s_bilan,
        s_difficultes, s_conclusion, s_questions,
    ]
    for fn in slides_fns:
        fn(prs)
    prs.save(pptx_file)
    size_kb = os.path.getsize(pptx_file) // 1024
    print(f'  OK  ->  {os.path.basename(pptx_file)}  ({size_kb} Ko)  [{len(slides_fns)} slides]')


# ═══════════════════════════════════════════════════════════════════════════════
#  PRÉSENTATION PDF PAYSAGE
# ═══════════════════════════════════════════════════════════════════════════════
def build_pres_pdf(pdf_file):
    from reportlab.pdfgen import canvas as rl_canvas
    PW, PH = landscape(A4)
    print(f'  Generation PDF paysage : {pdf_file}')
    c = rl_canvas.Canvas(pdf_file, pagesize=landscape(A4))
    slides_data = [
        ('ACADEMIQ — Rapport de Projet de Fin d\'Etudes', [
            ('Etudiant',    'MARIKO Lamine  |  CI0223065023'),
            ('Filiere',     'Master 1 — Genie Informatique'),
            ('Encadreur',   'Dr. ZEZE'),
            ('Institution', 'Universite Nangui Abrogoua  |  Juin 2026'),
            ('Stack',       'Django 4.x  .  MariaDB/InnoDB  .  Bootstrap 5  .  ReportLab'),
            ('Chiffres',    '24 tables  .  59 regles  .  7 roles  .  17 modules  .  840 resultats'),
        ]),
        ('Sommaire', [
            ('I.',   'Introduction et Contexte'),
            ('II.',  'Analyse et Specification des Besoins'),
            ('III.', 'Conception du Systeme (Methode Merise)'),
            ('IV.',  'Implementation et Developpement (Django)'),
            ('V.',   'Tests, Validation et Resultats'),
            ('VI.',  'Conclusion et Perspectives'),
        ]),
        ('Les 7 Roles Utilisateurs', [
            ('DIRECTION',      'Acces total : comptes, annees scolaires, historique'),
            ('ADMINISTRATION', 'Cours, emplois du temps, salles, matieres'),
            ('SCOLARITE',      'Inscriptions, absences, bulletins'),
            ('FINANCES',       'Tarifs, paiements, export bilan PDF'),
            ('ENSEIGNANT',     'Notes et absences pour ses cours uniquement'),
            ('ELEVE',          'Consultation de son propre dossier'),
            ('PARENT',         'Consultation du dossier de ses enfants'),
        ]),
        ('Bilan Quantitatif Final', [
            ('Tables BDD',      '24 tables, 95 attributs, 38 cles etrangeres'),
            ('Contraintes',     '12 UNIQUE composites, 1 CHECK [0;20], InnoDB ACID'),
            ('Regles metier',   '59 regles de gestion formalisees'),
            ('Roles',           '7 groupes Django avec perimetre cloisonne'),
            ('Modules livres',  '17 modules fonctionnels completement operationnels'),
            ('Donnees demo',    '840 RESULTAT_MATIERE avec rang calcule (0 NULL)'),
            ('Normalisation',   'Schema en 3FN verifie table par table'),
            ('Tests',           'Tests fonctionnels + integrite : tous conformes (OK)'),
        ]),
        ('Conclusion & Perspectives', [
            ('Bilan',        'Systeme complet et operationnel — 17 modules livres'),
            ('Securite',     'Controle acces 7 roles — aucun acces croise possible'),
            ('Automatisation','Moyennes, rangs, bulletins, notifications — signaux Django'),
            ('Court terme',  'Import CSV/Excel, tests unitaires, audit WCAG 2.1'),
            ('Moyen terme',  'API REST + app mobile (Django REST Framework)'),
            ('Long terme',   'Multi-etablissements, cloud, notifications push WebSocket'),
        ]),
    ]
    total = len(slides_data)
    for pn, (title, items) in enumerate(slides_data, 1):
        c.setFillColor(BF)
        c.rect(0, PH - 1.8*cm, PW, 1.8*cm, fill=1, stroke=0)
        c.setFillColor(colors.white)
        c.rect(0, 0, PW, PH - 1.8*cm, fill=1, stroke=0)
        c.setFillColor(BC)
        c.rect(0, PH - 1.83*cm, PW, 0.1*cm, fill=1, stroke=0)
        c.setFont('Helvetica-Bold', 16)
        c.setFillColor(colors.white)
        c.drawString(1.0*cm, PH - 1.3*cm, title)
        c.setFont('Helvetica', 8)
        c.setFillColor(colors.HexColor('#BDD7F5'))
        c.drawRightString(PW - 1.0*cm, PH - 1.3*cm, f'{pn} / {total}')
        y = PH - 2.6*cm
        for label, val in items:
            c.setFont('Helvetica-Bold', 10)
            c.setFillColor(BF)
            c.drawString(1.0*cm, y, f'{label}')
            c.setFont('Helvetica', 10)
            c.setFillColor(GT)
            c.drawString(5.5*cm, y, val)
            y -= 0.8*cm
        c.setStrokeColor(GR)
        c.setLineWidth(0.5)
        c.line(1.0*cm, 0.9*cm, PW - 1.0*cm, 0.9*cm)
        c.setFont('Helvetica', 7)
        c.setFillColor(BM)
        c.drawString(1.0*cm, 0.4*cm, 'ACADEMIQ  .  MARIKO Lamine  .  Juin 2026')
        c.drawRightString(PW - 1.0*cm, 0.4*cm, 'Universite Nangui Abrogoua  .  M1 GI')
        c.showPage()
    c.save()
    size_kb = os.path.getsize(pdf_file) // 1024
    print(f'  OK  ->  {os.path.basename(pdf_file)}  ({size_kb} Ko)')


# ═══════════════════════════════════════════════════════════════════════════════
#  POINT D'ENTREE
# ═══════════════════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    BASE = os.path.dirname(os.path.abspath(__file__))

    CDC_MD  = os.path.join(BASE, 'ACADEMIQ_V3_CahierDesCharges_BDD.md')
    RAP_MD  = os.path.join(BASE, 'ACADEMIQ_V3_Rapport_Final.md')
    CDC_PDF = os.path.join(BASE, 'ACADEMIQ_V3_CahierDesCharges_BDD.pdf')
    RAP_PDF = os.path.join(BASE, 'ACADEMIQ_V3_Rapport_Final.pdf')
    PPT_OUT = os.path.join(BASE, 'ACADEMIQ_V3_Presentation.pptx')
    PRS_PDF = os.path.join(BASE, 'ACADEMIQ_V3_Presentation_Slides.pdf')

    print('\n=== ACADEMIQ — Generation des documents ===\n')

    print('[1/4] Cahier des Charges PDF...')
    build_pdf(CDC_MD, CDC_PDF,
        doc_title='ACADEMIQ — Cahier des Charges & Conception BDD',
        right_header='MARIKO Lamine  .  CI0223065023',
        cover_fn=cover_cdc)

    print('[2/4] Rapport de Memoire PDF...')
    build_pdf(RAP_MD, RAP_PDF,
        doc_title='ACADEMIQ — Rapport de Projet de Fin d\'Etudes',
        right_header='MARIKO Lamine  .  CI0223065023',
        cover_fn=cover_rapport)

    print('[3/4] Presentation PPTX...')
    build_pptx(PPT_OUT)

    print('[4/4] Presentation PDF paysage...')
    build_pres_pdf(PRS_PDF)

    print('\n=== Tous les documents generes avec succes ===')
    for f in [CDC_PDF, RAP_PDF, PPT_OUT, PRS_PDF]:
        if os.path.exists(f):
            print(f'  {os.path.basename(f):52s} {os.path.getsize(f)//1024:>5} Ko')
